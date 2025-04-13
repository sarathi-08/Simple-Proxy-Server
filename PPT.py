from pptx import Presentation
from pptx.util import Inches, Pt

# Create a presentation object
prs = Presentation()

# Define content for slides
slides_content = [
    {
        "title": "Introduction to VLC",
        "content": [
            "Visible Light Communication (VLC) is a communication technology that utilizes visible light for data transmission.",
            "It employs Light Emitting Diodes (LEDs) for sending signals and photodetectors for receiving data.",
            "VLC offers a promising alternative in environments where radio frequency communication is limited.",
            "As a complementary technology, it can enhance existing wireless communication systems."
        ],
        "image": "vlc_overview_diagram.png"
    },
    {
        "title": "Working Principle of VLC",
        "content": [
            "VLC encodes data into light waves by modulating the intensity of LED lights at high speeds.",
            "This process is invisible to the naked eye, ensuring seamless communication.",
            "Photodetectors receive these light signals and decode the information back into data.",
            "The combination of modulation and photodetection allows for high-speed data transfer."
        ],
        "image": "vlc_working_principle.jpg"
    },
    {
        "title": "Advantages of VLC",
        "content": [
            "VLC offers several advantages, including increased bandwidth and security.",
            "Unlike traditional RF communication, VLC is not susceptible to interference from electronic devices.",
            "It provides a safe means of communication in environments like hospitals and airplanes.",
            "Additionally, the widespread use of LEDs makes VLC easily accessible."
        ],
        "image": None
    },
    {
        "title": "Applications of VLC",
        "content": [
            "VLC has diverse applications in various sectors, including healthcare, automotive, and smart cities.",
            "In healthcare, VLC can enable communication between devices in a secure manner.",
            "Smart cities can utilize VLC for lighting and data communication simultaneously.",
            "Moreover, it can be integrated into vehicle-to-vehicle communication systems."
        ],
        "image": None
    },
    {
        "title": "VLC vs. Traditional Wireless Communication",
        "content": [
            "VLC operates in the visible light spectrum, while traditional wireless communication relies on radio waves.",
            "This key difference allows VLC to offer higher data rates and improved security.",
            "VLC can also work in environments where RF communication is impractical, such as underwater or in-flight.",
            "However, VLC is limited by its line-of-sight requirement."
        ],
        "image": None
    },
    {
        "title": "Challenges in VLC Implementation",
        "content": [
            "Despite its advantages, VLC faces challenges like ambient light interference and limited range.",
            "The line-of-sight requirement can hinder communication in obstructed environments.",
            "Additionally, the technology requires advancements in photodetection for optimal performance.",
            "Understanding these challenges is essential for effective implementation."
        ],
        "image": None
    },
    {
        "title": "Future Trends in VLC",
        "content": [
            "The future of VLC looks promising, with ongoing research aimed at enhancing its capabilities.",
            "Developments in modulation techniques and photodetector technology will likely drive improvements.",
            "Integrating VLC with existing wireless technologies can lead to hybrid communication systems.",
            "As smart devices become more prevalent, VLC may play a crucial role in data communication."
        ],
        "image": None
    },
    {
        "title": "Case Studies",
        "content": [
            "Several case studies highlight the successful implementation of VLC in various settings.",
            "For instance, VLC has been used in museums for interactive exhibits, enhancing visitor experience.",
            "Another case study involves VLC applications in smart lighting systems for urban environments.",
            "These examples showcase the versatility and potential of VLC technology."
        ],
        "image": None
    },
    {
        "title": "Conclusion",
        "content": [
            "In conclusion, VLC presents a viable solution for modern communication needs.",
            "Its unique advantages and diverse applications make it a valuable addition to wireless technologies.",
            "As research progresses, VLC's role in communication will continue to expand.",
            "The integration of VLC with existing technologies can lead to innovative solutions."
        ],
        "image": None
    },
    {
        "title": "References",
        "content": [
            "1. Author A, Title of Book, Year.",
            "2. Author B, 'Title of Article,' Journal Name, Year.",
            "3. Author C, Title of Website, Date Accessed.",
            "4. Author D, Title of Conference Paper, Year."
        ],
        "image": None
    },
    {
        "title": "Technical Specifications",
        "content": [
            "Key specifications of VLC include data rate, transmission range, and modulation techniques.",
            "Common standards for VLC technology involve IEEE 802.15.7 and visible light communication protocols.",
            "These specifications ensure compatibility and performance across different devices.",
            "Understanding these technical aspects is crucial for effective VLC deployment."
        ],
        "image": "vlc_technical_specifications.png"
    },
    {
        "title": "Potential Impact of VLC on Society",
        "content": [
            "VLC has the potential to revolutionize communication by providing high-speed data transfer in various environments.",
            "Its applications in smart cities can lead to improved connectivity and efficiency.",
            "Moreover, VLC can enhance safety in critical environments like hospitals by minimizing RF interference.",
            "As VLC technology matures, its societal impact will continue to grow."
        ],
        "image": None
    }
]

# Adding slides to the presentation
for slide_info in slides_content:
    slide_layout = prs.slide_layouts[5]  # Using a blank slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title = slide.shapes.title
    title.text = slide_info["title"]
    
    # Add content
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(4))
    text_frame = textbox.text_frame
    for line in slide_info["content"]:
        p = text_frame.add_paragraph()
        p.text = line
        p.space_after = Pt(14)  # Space after each paragraph

    # Add image if specified
    if slide_info["image"]:
        slide.shapes.add_picture(slide_info["image"], Inches(1), Inches(2), width=Inches(6), height=Inches(3))

# Save the presentation
prs.save("VLC_Presentation.pptx")
